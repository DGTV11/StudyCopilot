import gradio as gr
import ollama, os
from pptx import Presentation

fh_is_regenerating = False

def gen_flashcards(filepaths, notes):
	global fh_is_regenerating
	if fh_is_regenerating:
		gr.Error('Please wait until flashcards_helper model is ready.')

	pptx_files_notes = ''
	if filepaths:
		len_filepaths = len(filepaths)
		for i, filepath in enumerate(filepaths, start=1):
			pptx_files_notes += os.path.basename(filepath).split('.')[0] + '\n'
			slides = Presentation(filepath).slides
			gr.Info(f'Grabbing text from slideshow {i}/{len_filepaths} ({len(slides)} slides)...')
			for slide in slides:
				for shape in slide.shapes:
					if hasattr(shape, 'text'):
						pptx_files_notes += shape.text + '\n'

	gr.Info('Reading slide texts and notes...')
	stream = ollama.chat(
			model='flashcards_helper',
			messages=[{'role': 'user', 'content': f'Text: {pptx_files_notes + notes}\n\nA deck of flashcards:'}],
			stream=True,
	)
	res_stream = ''

	gr.Info('Generating flashcards...')
	for chunk in stream:
		res_stream += chunk['message']['content']
		yield res_stream

def regen_flashcards_helper():
	global fh_is_regenerating

	fh_is_regenerating = True
	gr.Info('Regenerating flashcards_helper from Modelfile...')

	with open(os.path.join(os.path.dirname(__file__), 'Modelfile')) as file:
		ollama.create(model='flashcards_helper', modelfile=file.read())

	fh_is_regenerating = False
	gr.Info('Finished regenerating!')

if __name__ == '__main__':
	print('Checking if flashcards_helper model exists...')
	if 'flashcards_helper:latest' not in [model['name'] for model in ollama.list()['models']]:
		with open(os.path.join(os.path.dirname(__file__), 'Modelfile')) as file:
			print('Creating flashcards_helper from Modelfile...')
			ollama.create(model='flashcards_helper', modelfile=file.read())	

	with gr.Blocks(theme=gr.themes.Default(primary_hue="red")) as studyCopilot:
		gr.Markdown(
		"""
		# Study Copilot
		""")

		with gr.Tab('Flashcard Generator'):
			inp_cards_slides = gr.File(label="Slides:", file_count='multiple', file_types=['.pptx'])	
			inp_cards_words = gr.Textbox(label="Notes:")

			with gr.Row():
				clear_inputs_btn = gr.ClearButton(components=[inp_cards_slides, inp_cards_words])
				generate_cards_btn = gr.Button(value="Generate Flashcards")
			
			stop = gr.Button(value="Stop", variant='primary')
			out_cards = gr.Markdown(label="A deck of flashcards:")
			clear_cards_btn = gr.ClearButton(components=[out_cards])
			regen_flashcards_helper_btn = gr.Button(value="Regenerate flashcards_helper", variant='secondary')

		regen_flashcards_helper_btn.click(fn=regen_flashcards_helper, inputs=[], outputs=[])
		gen_cards_event = generate_cards_btn.click(fn=gen_flashcards, inputs=[inp_cards_slides, inp_cards_words], outputs=[out_cards])
		stop.click(fn=None, inputs=None, outputs=None, cancels=[gen_cards_event])

	print('Launching...')
	studyCopilot.launch()