import gradio as gr
import ollama, os
from pptx import Presentation

def gen_flashcards(slides_filepaths, notes):
	global fh_is_regenerating
	if fh_is_regenerating:
		raise gr.Error('Please wait until flashcards_helper model is ready.')
	elif (not slides_filepaths) and (not notes):
		raise gr.Error('Please provide at least one file and/or text.')
	else:
		slides_notes = ''

		if slides_filepaths:
			len_filepaths = len(slides_filepaths)
			for i, filepath in enumerate(slides_filepaths, start=1):
				slides_notes += os.path.basename(filepath).split('.')[0] + '\n'
				slides = Presentation(filepath).slides
				gr.Info(f'Grabbing text from slideshow {i}/{len_filepaths} ({len(slides)} slides)...')
				for slide in slides:
					for shape in slide.shapes:
						if hasattr(shape, 'text'):
							slides_notes += shape.text + '\n'

		gr.Info('Sending files and text to flashcards_helper...')
		stream = ollama.chat(
				model='flashcards_helper',
				messages=[{'role': 'user', 'content': f'Text: {slides_notes + notes}\n\nA deck of flashcards:'}],
				stream=True,
		)
		res_stream = ''

		gr.Info('Generating flashcards...')
		for chunk in stream:
			res_stream += chunk['message']['content']
			yield res_stream