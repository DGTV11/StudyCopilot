import os
from time import time

import gradio as gr
from pptx import Presentation
from semantic_text_splitter import TextSplitter
# from tokenizers import Tokenizer
from transformers import AutoTokenizer

import modules.logging as log
from modules.host import HOST
from modules.config import CONFIG

SYSTEM_PROMPT = """I want you to create a deck of flashcards from the text.
The text can be very long and in different formats such as paragraphs, markdown, bulletpoints or a combination of them.

Instructions to create a deck of flashcards:
- Write your thought process before creating the flashcards.
- Keep the flashcards simple, clear, and focused on the most important information.
- Make sure the questions are specific and unambiguous.
- Use simple and direct language to make the cards easy to read and understand.
- Answers should contain only a single key fact/name/concept/term.
- Make only one deck of flashcards; do not make extra sets of texts and decks of flashcards on your own.

Let's do it step by step when creating a deck of flashcards:
1. Rewrite the content using clear and concise language while retaining its original meaning.
2. Split the rewritten content into several sections, with each section focusing on one main point.
3. Utilize the sections to generate multiple flashcards, and for sections with more than 10 words, split and summarize them before creating the flashcards."""

EXAMPLE_PROMPT = "The characteristics of the Dead Sea: Salt lake located on the border between Israel and Jordan. Its shoreline is the lowest point on the Earth's surface, averaging 396 m below sea level. It is 74 km long. It is seven times as salty (30% by volume) as the ocean. Its density keeps swimmers afloat. Only simple organisms can live in its saline waters."

EXAMPLE_RESPONSE = """# THOUGHT PROCESS
To create the flashcards for the characteristics of the Dead Sea, I will follow these steps:

1. **Rewrite the content using clear and concise language:**
   - The Dead Sea is a salt lake located between Israel and Jordan.
   - Its shoreline is the lowest point on Earth's surface, averaging 396 meters below sea level.
   - The Dead Sea is 74 kilometers long.
   - It is seven times saltier than the ocean, with a salt concentration of 30% by volume.
   - The high density of the water keeps swimmers afloat.
   - Only simple organisms can survive in its saline waters.

2. **Split the rewritten content into several sections:**
   - Location and geographical context
   - Shoreline characteristics
   - Length of the Dead Sea
   - Salinity compared to the ocean
   - Density and buoyancy
   - Organisms living in the Dead Sea

3. **Utilize the sections to generate flashcards:**

   - **Location and geographical context:**
     - Where is the Dead Sea located? (Answer: on the border between Israel and Jordan)
     - What is the lowest point on the Earth's surface? (Answer: The Dead Sea shoreline)

   - **Shoreline characteristics:**
     - What is the average level on which the Dead Sea is located? (Answer: 396 meters below sea level)

   - **Length of the Dead Sea:**
     - How long is the Dead Sea? (Answer: 74 km)

   - **Salinity compared to the ocean:**
     - How much saltier is the Dead Sea compared with the oceans? (Answer: 7 times)
     - What is the volume content of salt in the Dead Sea? (Answer: 30%)

   - **Density and buoyancy:**
     - Why can the Dead Sea keep swimmers afloat? (Answer: due to high salt content)

   - **Organisms living in the Dead Sea:**
     - Why is the Dead Sea called Dead? (Answer: because only simple organisms can live in it)
     - Why only simple organisms can live in the Dead Sea? (Answer: because of high salt content)

I will now generate the deck of flashcards based on the generated thought process.

---

# A deck of flashcards
|Question|Answer|
|---|---|
|Where is the Dead Sea located?|on the border between Israel and Jordan|
|What is the lowest point on the Earth's surface?|The Dead Sea shoreline|
|What is the average level on which the Dead Sea is located?|396 meters (below sea level)|
|How long is the Dead Sea?|74 km|
|How much saltier is the Dead Sea as compared with the oceans?|7 times|
|What is the volume content of salt in the Dead Sea?|30%|
|Why can the Dead Sea keep swimmers afloat?|due to high salt content|
|Why is the Dead Sea called Dead?|because only simple organisms can live in it|
|Why only simple organisms can live in the Dead Sea?|because of high salt content|
"""

INITIAL_CHAT_HISTORY = [
    {"role": "system", "content": SYSTEM_PROMPT},
    {"role": "user", "content": EXAMPLE_PROMPT},
    {"role": "assistant", "content": EXAMPLE_RESPONSE}
]

def show_markdown(text):
    return text if text.startswith("```") and text.endswith("```") else "```" + text


def send_to_model(flashcards_helper_model, ctx_window, text):
    stream = HOST.chat(
        model=flashcards_helper_model,
        messages=INITIAL_CHAT_HISTORY + [{"role": "user", "content": text}],
        stream=True,
        options={"num_ctx": ctx_window},
    )
    res_stream = ""

    log.log_info("Flashcard Generator", "Generating flashcards...")
    for chunk in stream:
        yield chunk["message"]["content"]


def gen_flashcards(
    flashcards_helper_model, slides_filepaths, image_filepaths, notes, get_slides_images
):
    gen_start_time = time()
    match flashcards_helper_model:
        case "openhermes":
            tokenizer = AutoTokenizer.from_pretrained(
                "teknium/OpenHermes-2.5-Mistral-7B",
                atoken=CONFIG["huggingface_user_access_token"],
            )._tokenizer
            ctx_window = 4096
        case "phi3":
            tokenizer = AutoTokenizer.from_pretrained(
                "microsoft/Phi-3-mini-4k-instruct",
                token=CONFIG["huggingface_user_access_token"],
            )._tokenizer
            ctx_window = 4096
        case _:
            raise gr.Error(f"{flashcards_helper_model} is not a supported model.")

    num_token_func = lambda text: len(tokenizer.encode(text).ids)

    if (not slides_filepaths) and (not notes) and (not image_filepaths):
        raise gr.Error(
            "Please provide at least one slideshow and/or image file and/or text."
        )
    else:
        res_stream = ""
        if slides_filepaths:
            len_filepaths = len(slides_filepaths)
            for i, filepath in enumerate(slides_filepaths, start=1):
                slides_name = os.path.basename(filepath).split(".")[0] + "\n"
                slides = Presentation(filepath).slides
                len_slides = len(slides)
                log.log_info(
                    "Flashcard Generator",
                    f"Reading slideshow {i}/{len_filepaths} ({len(slides)} slides)",
                )
                slides_notes = slides_name
                slides_notes_start_slide = 1
                for j, slide in enumerate(slides, start=1):
                    log.log_info(
                        "Flashcard Generator",
                        f"Reading slide {j}/{len_slides}...",
                        debug_only=True,
                    )
                    slide_notes = ""

                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_notes += shape.text + "\n"
                        if get_slides_images and hasattr(shape, "image"):
                            log.log_info(
                                "Flashcard Generator",
                                f"Grabbing and describing image from slide {j}/{len_slides}",
                            )

                            start_time = time()
                            slide_notes += (
                                "\n"
                                + HOST.generate(
                                    model="llava-llama3",
                                    prompt="Describe the content of this image in detail as if it were text",
                                    images=[shape.image.blob],
                                )["response"]
                                + "\n"
                            )
                            end_time = time()

                            log.log_info(
                                "Flashcard Generator",
                                f"Finished describing image! ({round(end_time-start_time, 2)}s)",
                            )

                    slide_notes_num_tokens = num_token_func(slide_notes)
                    slides_notes_num_tokens = num_token_func(slides_notes)

                    if slides_notes_num_tokens + slide_notes_num_tokens > ctx_window:
                        log.log_info(
                            "Flashcard Generator",
                            f"Sending slides {slides_notes_start_slide}-{j-1} of slideshow {i}/{len_filepaths} to flashcards_helper_{flashcards_helper_model} ({slides_notes_num_tokens}/{ctx_window} tokens)",
                        )

                        start_time = time()
                        for res in send_to_model(
                            flashcards_helper_model, ctx_window, slides_notes
                        ):
                            res_stream += res
                            yield res_stream
                        end_time = time()

                        log.log_info(
                            "Flashcard Generator",
                            f"Finished generating batch of flashcards ({round(end_time-start_time, 2)}s)",
                        )

                        res_stream += "\n\n"
                        slides_notes = slides_name + slide_notes
                        slides_notes_start_slide = j
                    else:
                        slides_notes += slide_notes

                    if j == len_slides:
                        log.log_info(
                            "Flashcard Generator",
                            f"Sending slides {slides_notes_start_slide}-{j} of slideshow {i}/{len_filepaths} to flashcards_helper_{flashcards_helper_model} ({slides_notes_num_tokens + slide_notes_num_tokens}/{ctx_window} tokens)",
                        )

                        start_time = time()
                        for res in send_to_model(
                            flashcards_helper_model, ctx_window, slides_notes
                        ):
                            res_stream += res
                            yield res_stream
                        end_time = time()

                        log.log_info(
                            "Flashcard Generator",
                            f"Finished generating batch of flashcards ({round(end_time-start_time, 2)}s)",
                        )

                        res_stream += "\n\n"

        if image_filepaths:
            len_filepaths = len(image_filepaths)
            images_notes_start_image = 1
            images_notes = ""
            for i, filepath in enumerate(image_filepaths, start=1):
                image_notes = os.path.basename(filepath).split(".")[0] + "\n"

                log.log_info(
                    "Flashcard Generator",
                    f"Grabbing and describing image {i}/{len_filepaths}...",
                )

                start_time = time()
                image_notes = (
                    "\n"
                    + HOST.generate(
                        model="llava-llama3",
                        prompt="Describe the content of this image in detail as if it were text",
                        images=[filepath],
                    )["response"]
                    + "\n"
                )
                end_time = time()

                log.log_info(
                    "Flashcard Generator",
                    f"Finished describing image! ({round(end_time-start_time, 2)}s)",
                )

                image_notes_num_tokens = num_token_func(image_notes)
                images_notes_num_tokens = num_token_func(images_notes)

                if images_notes_num_tokens + image_notes_num_tokens > ctx_window:
                    log.log_info(
                        "Flashcard Generator",
                        f"Sending images {images_notes_start_image}/{len_filepaths}-{i-1}/{len_filepaths} to flashcards_helper_{flashcards_helper_model} ({images_notes_num_tokens}/{ctx_window} tokens)",
                    )

                    start_time = time()
                    for res in send_to_model(
                        flashcards_helper_model, ctx_window, images_notes
                    ):
                        res_stream += res
                        yield res_stream
                    end_time = time()

                    log.log_info(
                        "Flashcard Generator",
                        f"Finished generating batch of flashcards ({round(end_time-start_time, 2)}s)",
                    )

                    res_stream += "\n\n"
                    images_notes = image_notes
                    images_notes_start_image = i
                else:
                    images_notes += image_notes

                if i == len_filepaths:
                    log.log_info(
                        "Flashcard Generator",
                        f"Sending images {images_notes_start_image}/{len_filepaths}-{i}/{len_filepaths} to flashcards_helper_{flashcards_helper_model} ({images_notes_num_tokens + image_notes_num_tokens}/{ctx_window} tokens)",
                    )

                    start_time = time()
                    for res in send_to_model(
                        flashcards_helper_model, ctx_window, images_notes
                    ):
                        res_stream += res
                        yield res_stream
                    end_time = time()

                    log.log_info(
                        "Flashcard Generator",
                        f"Finished generating batch of flashcards ({round(end_time-start_time, 2)}s)",
                    )

                    res_stream += "\n\n"

        if notes.strip():
            splitter = TextSplitter.from_huggingface_tokenizer(tokenizer, ctx_window)

            chunks = splitter.chunks(notes)
            len_chunks = len(chunks)

            for i, chunk in enumerate(chunks, start=1):
                chunk_num_tokens = num_token_func(chunk)
                log.log_info(
                    "Flashcard Generator",
                    f"Sending chunk {i}/{len_chunks} of textual notes to flashcards_helper_{flashcards_helper_model} ({chunk_num_tokens}/{ctx_window} tokens)",
                )

                start_time = time()
                for res in send_to_model(flashcards_helper_model, ctx_window, chunk):
                    res_stream += res
                    yield res_stream
                end_time = time()

                log.log_info(
                    "Flashcard Generator",
                    f"Finished generating batch of flashcards ({round(end_time-start_time, 2)}s)",
                )

                res_stream += "\n\n"

        gen_end_time = time()
        log.log_info(
            "Flashcard Generator",
            f"Finished generating flashcards ({round(gen_end_time-gen_start_time, 2)}s)",
        )
        yield res_stream
